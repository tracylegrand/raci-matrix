import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# Initialize session state
if 'raci_data' not in st.session_state:
    st.session_state.raci_data = pd.DataFrame()
if 'functions' not in st.session_state:
    st.session_state.functions = []
if 'stakeholders' not in st.session_state:
    st.session_state.stakeholders = []
if 'function_input_key' not in st.session_state:
    st.session_state.function_input_key = 0
if 'stakeholder_input_key' not in st.session_state:
    st.session_state.stakeholder_input_key = 0
if 'refocus_function' not in st.session_state:
    st.session_state.refocus_function = False
if 'refocus_stakeholder' not in st.session_state:
    st.session_state.refocus_stakeholder = False
if 'last_raci_data_hash' not in st.session_state:
    st.session_state.last_raci_data_hash = None
if 'brian_joke' not in st.session_state:
    st.session_state.brian_joke = False

# RACI options - keys are what get stored, values are what display in dropdown
RACI_OPTIONS = {
    '': '',
    'R': 'R - Responsible',
    'A': 'A - Accountable',
    'C': 'C - Consulted',
    'I': 'I - Informed'
}

# RACI labels for display (maps letter to full label)
RACI_LABELS = {
    'R': 'R - Responsible',
    'A': 'A - Accountable',
    'C': 'C - Consulted',
    'I': 'I - Informed'
}

# Color scheme for RACI
RACI_COLORS = {
    'R': '#FFE5B4',  # Light orange
    'A': '#FFB6C1',  # Light pink
    'C': '#B0E0E6',  # Powder blue
    'I': '#E0E0E0'   # Light gray
}

def create_raci_matrix(functions, stakeholders):
    """Create an empty RACI matrix DataFrame"""
    if not functions or not stakeholders:
        return pd.DataFrame()
    
    # Create matrix with functions as rows and stakeholders as columns
    matrix = pd.DataFrame(
        index=functions,
        columns=stakeholders,
        data=''
    )
    return matrix

def validate_raci_row(row_data):
    """Validate that a row has at most one 'A' (Accountable) role"""
    accountable_count = sum(1 for val in row_data if str(val).strip() == 'A')
    return accountable_count <= 1

def validate_raci_matrix(df):
    """Validate the entire RACI matrix for correctness"""
    errors = []
    for function in df.index:
        row_data = df.loc[function].values
        # Extract letter from value (could be 'A' or 'A - Accountable')
        accountable_count = 0
        for val in row_data:
            val_str = str(val).strip()
            # Check if it starts with 'A' (could be 'A' or 'A - Accountable')
            if val_str and (val_str == 'A' or val_str.startswith('A -')):
                accountable_count += 1
        if accountable_count > 1:
            errors.append(f"Function '{function}' has {accountable_count} Accountable stakeholders. Only 1 is allowed.")
    return errors

def parse_raci_value(value):
    """Parse RACI value from imported file - handles both formats"""
    if pd.isna(value) or value == '':
        return ''
    
    val_str = str(value).strip().upper()
    
    # If it's just a single letter, return it
    if len(val_str) == 1 and val_str in ['R', 'A', 'C', 'I']:
        return RACI_LABELS.get(val_str, val_str)
    
    # If it starts with a RACI letter, extract it
    if val_str and val_str[0] in ['R', 'A', 'C', 'I']:
        letter = val_str[0]
        return RACI_LABELS.get(letter, letter)
    
    # Try to match full labels
    val_lower = val_str.lower()
    if 'responsible' in val_lower:
        return RACI_LABELS['R']
    elif 'accountable' in val_lower:
        return RACI_LABELS['A']
    elif 'consulted' in val_lower:
        return RACI_LABELS['C']
    elif 'informed' in val_lower:
        return RACI_LABELS['I']
    
    return ''

def import_from_spreadsheet(uploaded_file):
    """Import RACI matrix from Excel or CSV file"""
    try:
        # Reset file pointer to beginning
        uploaded_file.seek(0)
        
        # Determine file type and read accordingly
        if uploaded_file.name.endswith('.csv'):
            df = pd.read_csv(uploaded_file, index_col=0)
        elif uploaded_file.name.endswith(('.xlsx', '.xls')):
            # Try to read the first sheet
            uploaded_file.seek(0)  # Reset again for Excel
            df = pd.read_excel(uploaded_file, index_col=0, sheet_name=0)
        else:
            return False, "Unsupported file format. Please use Excel (.xlsx, .xls) or CSV (.csv) files."
        
        # Check if DataFrame is empty
        if df.empty:
            return False, "The uploaded file appears to be empty."
        
        # Drop rows where the index is NaN or empty
        df = df[df.index.notna()]
        df = df[df.index.astype(str).str.strip() != '']
        
        # Filter out legend rows and unnamed rows
        # Legend rows typically have index starting with "Legend" or contain "=" or are in common legend formats
        legend_patterns = ['legend', 'r =', 'a =', 'c =', 'i =', 'responsible', 'accountable', 'consulted', 'informed']
        
        def is_legend_row(idx):
            """Check if a row index looks like a legend row"""
            idx_str = str(idx).strip().lower()
            # Check if index contains legend patterns
            if any(pattern in idx_str for pattern in legend_patterns):
                return True
            # Check if index contains "=" (common in legends like "R = Responsible")
            if '=' in idx_str:
                return True
            return False
        
        def is_unnamed_row(idx):
            """Check if a row index is an unnamed/empty row"""
            idx_str = str(idx).strip()
            # Check for "Unnamed" pattern (case-insensitive) - handles "Unnamed: 3", "Unnamed: 4", etc.
            if 'unnamed' in idx_str.lower():
                return True
            # Check for empty or NaN values
            if not idx_str or idx_str.lower() in ['nan', 'none', '']:
                return True
            return False
        
        # Filter out legend rows and unnamed rows from the DataFrame
        df_filtered = df[~df.index.map(lambda idx: is_legend_row(idx) or is_unnamed_row(idx))]
        
        # Also check if row has mostly empty cells or legend-like content in cells
        # (Sometimes legend is in a row with mixed content)
        def row_looks_like_legend(row):
            """Check if a row looks like a legend row based on cell content"""
            row_str = ' '.join([str(val).strip().lower() for val in row.values if pd.notna(val)])
            # If row contains legend patterns, it's likely a legend
            if any(pattern in row_str for pattern in ['r =', 'a =', 'c =', 'i =', 'responsible', 'accountable', 'consulted', 'informed']):
                return True
            return False
        
        def row_is_empty_or_unnamed(row, row_idx):
            """Check if a row is mostly empty or has no meaningful data"""
            # First check if the index itself is unnamed
            if is_unnamed_row(row_idx):
                return True
            
            # Count non-empty, non-NaN values
            non_empty_values = [val for val in row.values if pd.notna(val) and str(val).strip() != '']
            non_empty_count = len(non_empty_values)
            
            # If row has no non-empty values, it's empty
            if non_empty_count == 0:
                return True
            
            # Check if all values in the row are empty strings or NaN
            all_empty = all(pd.isna(val) or str(val).strip() == '' for val in row.values)
            if all_empty:
                return True
            
            # Check if the row index is empty/NaN but row has some data (might be a data row)
            # But if index is unnamed and row has minimal data, it's likely still an unwanted row
            idx_str = str(row_idx).strip()
            if 'unnamed' in idx_str.lower() and non_empty_count <= 1:
                return True
            
            return False
        
        # Additional filter: remove rows that look like legends or are empty/unnamed
        mask = ~df_filtered.apply(lambda row: row_looks_like_legend(row) or row_is_empty_or_unnamed(row, row.name), axis=1)
        df_filtered = df_filtered[mask]
        
        if df_filtered.empty:
            return False, "No valid data rows found after filtering legend. Please ensure your file contains function names in the first column."
        
        # Extract functions from index (first column) - now filtered
        functions = [str(idx).strip() for idx in df_filtered.index if str(idx).strip() and str(idx).strip().lower() not in ['nan', 'none', '']]
        
        # Extract stakeholders from column headers - filter out blank/empty/unnamed columns
        def is_valid_stakeholder(col):
            """Check if a column header is a valid stakeholder name"""
            col_str = str(col).strip()
            # Check for empty or NaN
            if not col_str or col_str.lower() in ['nan', 'none', '']:
                return False
            # Check for "Unnamed" pattern
            if 'unnamed' in col_str.lower():
                return False
            return True
        
        # Filter columns by header name (removes blank/unnamed columns)
        valid_cols = [col for col in df_filtered.columns if is_valid_stakeholder(col)]
        df_filtered = df_filtered[valid_cols]
        
        # Also filter out columns that are completely empty (all NaN or empty values)
        # This catches columns with valid headers but no data
        cols_with_data = []
        for col in df_filtered.columns:
            col_data = df_filtered[col]
            # Check if column has any non-empty, non-NaN values
            has_values = any(pd.notna(val) and str(val).strip() != '' for val in col_data.values)
            if has_values:
                cols_with_data.append(col)
        
        # Use columns with data (or at least valid headers)
        df_filtered = df_filtered[cols_with_data] if cols_with_data else df_filtered[valid_cols]
        
        # Extract stakeholders from remaining valid columns
        stakeholders = [str(col).strip() for col in df_filtered.columns if is_valid_stakeholder(col)]
        
        if not functions:
            return False, "No functions found in the file. Please ensure the first column contains function names."
        
        if not stakeholders:
            return False, "No stakeholders found in the file. Please ensure the first row contains stakeholder names."
        
        # Parse RACI values - convert to full label format
        raci_df = df_filtered.copy()
        for col in raci_df.columns:
            for idx in raci_df.index:
                original_value = raci_df.loc[idx, col]
                parsed_value = parse_raci_value(original_value)
                raci_df.loc[idx, col] = parsed_value
        
        # Update session state
        st.session_state.functions = functions
        st.session_state.stakeholders = stakeholders
        st.session_state.raci_data = raci_df.fillna('')
        
        return True, f"Successfully imported {len(functions)} functions and {len(stakeholders)} stakeholders!"
        
    except Exception as e:
        return False, f"Error importing file: {str(e)}"

def export_to_excel(df, filename='raci_matrix.xlsx'):
    """Export RACI matrix to Excel with formatting"""
    if df.empty:
        raise ValueError("Cannot export empty matrix. Please add functions and stakeholders first.")
    
    output = BytesIO()
    
    try:
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='RACI Matrix', index=True)
            worksheet = writer.sheets['RACI Matrix']
            
            # Set column widths
            worksheet.column_dimensions['A'].width = 25
            for col in range(2, len(df.columns) + 2):
                worksheet.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 15
            
            # Style header row
            header_fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
            header_font = Font(color='FFFFFF', bold=True, size=11)
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Format header
            for cell in worksheet[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = border
            
            # Format index column
            index_fill = PatternFill(start_color='D9E1F2', end_color='D9E1F2', fill_type='solid')
            index_font = Font(bold=True, size=11)
            
            for row in range(2, len(df) + 2):
                cell = worksheet[f'A{row}']
                cell.fill = index_fill
                cell.font = index_font
                cell.alignment = Alignment(horizontal='left', vertical='center')
                cell.border = border
            
            # Format data cells with RACI colors
            for row_idx, row in enumerate(df.index, start=2):
                for col_idx, col in enumerate(df.columns, start=2):
                    cell = worksheet.cell(row=row_idx, column=col_idx)
                    value = str(df.loc[row, col]).strip()
                    
                    # Extract letter and convert to label for display
                    raci_letter = ''
                    display_value = value
                    if value:
                        if value.startswith('R -') or value == 'R':
                            raci_letter = 'R'
                            display_value = 'R - Responsible'
                        elif value.startswith('A -') or value == 'A':
                            raci_letter = 'A'
                            display_value = 'A - Accountable'
                        elif value.startswith('C -') or value == 'C':
                            raci_letter = 'C'
                            display_value = 'C - Consulted'
                        elif value.startswith('I -') or value == 'I':
                            raci_letter = 'I'
                            display_value = 'I - Informed'
                        elif len(value) == 1 and value in ['R', 'A', 'C', 'I']:
                            raci_letter = value
                            display_value = RACI_LABELS.get(value, value)
                        else:
                            raci_letter = value[0] if len(value) > 0 else ''
                            if raci_letter in RACI_LABELS:
                                display_value = RACI_LABELS[raci_letter]
                    
                    cell.value = display_value
                    
                    if raci_letter in RACI_COLORS:
                        hex_color = RACI_COLORS[raci_letter].replace('#', '')
                        cell.fill = PatternFill(
                            start_color=hex_color,
                            end_color=hex_color,
                            fill_type='solid'
                        )
                    
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                    cell.border = border
                    cell.font = Font(size=11, bold=(raci_letter in ['R', 'A']))
            
            # Add legend
            legend_row = len(df) + 3
            worksheet.cell(row=legend_row, column=1, value='Legend:')
            worksheet.cell(row=legend_row, column=1).font = Font(bold=True, size=11)
            
            legend_items = ['R = Responsible', 'A = Accountable', 'C = Consulted', 'I = Informed']
            for idx, item in enumerate(legend_items, start=2):
                cell = worksheet.cell(row=legend_row, column=idx, value=item)
                cell.font = Font(size=10)
    except Exception as e:
        st.error(f"Error exporting to Excel: {str(e)}")
        raise
    
    output.seek(0)
    return output

def export_to_powerpoint(df, filename='raci_matrix.pptx'):
    """Export RACI matrix to PowerPoint presentation"""
    if df.empty:
        raise ValueError("Cannot export empty matrix. Please add functions and stakeholders first.")
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Use blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(9)
        title_height = Inches(0.5)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = "RACI Matrix"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Calculate table dimensions
        table_left = Inches(0.5)
        table_top = Inches(1)
        table_width = Inches(9)
        table_height = Inches(5)
        
        # Create table
        num_rows = len(df) + 1  # +1 for header
        num_cols = len(df.columns) + 1  # +1 for index column
        
        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height).table
        
        # Set column widths
        index_col_width = table_width / num_cols * 1.5
        data_col_width = (table_width - index_col_width) / (num_cols - 1)
        
        # Fill header row
        table.cell(0, 0).text = "Function"
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = RGBColor(54, 96, 146)
        
        for col_idx, stakeholder in enumerate(df.columns, start=1):
            cell = table.cell(0, col_idx)
            cell.text = stakeholder
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
        
        # Fill data rows
        for row_idx, function in enumerate(df.index, start=1):
            # Index column
            table.cell(row_idx, 0).text = function
            table.cell(row_idx, 0).fill.solid()
            table.cell(row_idx, 0).fill.fore_color.rgb = RGBColor(217, 225, 242)
            
            # Data cells
            for col_idx, stakeholder in enumerate(df.columns, start=1):
                cell = table.cell(row_idx, col_idx)
                value = str(df.loc[function, stakeholder]).strip()
                
                # Value could be just letter ('R') or full label ('R - Responsible')
                # Extract the letter for color matching
                raci_letter = ''
                display_text = value
                
                if value:
                    # Check if it's a full label format
                    if value.startswith('R -') or value == 'R - Responsible':
                        raci_letter = 'R'
                        display_text = 'R - Responsible'
                    elif value.startswith('A -') or value == 'A - Accountable':
                        raci_letter = 'A'
                        display_text = 'A - Accountable'
                    elif value.startswith('C -') or value == 'C - Consulted':
                        raci_letter = 'C'
                        display_text = 'C - Consulted'
                    elif value.startswith('I -') or value == 'I - Informed':
                        raci_letter = 'I'
                        display_text = 'I - Informed'
                    elif len(value) == 1 and value in ['R', 'A', 'C', 'I']:
                        # Just the letter - convert to full label
                        raci_letter = value
                        display_text = RACI_LABELS.get(value, value)
                    else:
                        # Try to extract first character
                        raci_letter = value[0] if len(value) > 0 else ''
                        if raci_letter in RACI_LABELS:
                            display_text = RACI_LABELS[raci_letter]
                
                cell.text = display_text
                
                # Apply color based on the letter
                if raci_letter in RACI_COLORS:
                    color = RACI_COLORS[raci_letter]
                    rgb = tuple(int(color[i:i+2], 16) for i in (1, 3, 5))
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*rgb)
        
        # Format all cells with larger font
        for row in range(num_rows):
            for col in range(num_cols):
                cell = table.cell(row, col)
                # Increased font size for data cells (was Pt(10), now Pt(14))
                # Header and index keep larger size
                if row == 0 or col == 0:
                    # Header and index cells
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.bold = True
                    if row == 0:
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                else:
                    # Data cells - larger font
                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                    cell.text_frame.paragraphs[0].font.bold = False
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add legend
        legend_top = Inches(6.5)
        legend_left = Inches(0.5)
        legend_width = Inches(9)
        legend_height = Inches(0.5)
        
        legend_box = slide.shapes.add_textbox(legend_left, legend_top, legend_width, legend_height)
        legend_frame = legend_box.text_frame
        legend_frame.text = "Legend: R = Responsible | A = Accountable | C = Consulted | I = Informed"
        legend_para = legend_frame.paragraphs[0]
        legend_para.font.size = Pt(9)
        legend_para.alignment = PP_ALIGN.CENTER
        
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error exporting to PowerPoint: {str(e)}")
        raise

# Streamlit UI
st.set_page_config(page_title="RACI Matrix Builder", page_icon="📊", layout="wide")

# Version and author info
VERSION = "1.1.0-joke"

# Custom CSS for version and author display in upper right
st.markdown(f"""
    <style>
    .version-info {{
        position: fixed !important;
        top: 10px !important;
        right: 20px !important;
        background-color: rgba(255, 255, 255, 0.95) !important;
        padding: 8px 12px !important;
        border-radius: 5px !important;
        font-size: 12px !important;
        color: #666 !important;
        z-index: 9999 !important;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1) !important;
        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif !important;
        white-space: nowrap !important;
        display: block !important;
        visibility: visible !important;
    }}
    </style>
    <div class="version-info">
        <strong>v{VERSION}</strong> | By: <a href="https://www.linkedin.com/in/tracylegrand/" target="_blank" style="color: #0066cc; text-decoration: underline;">TDLG</a>
    </div>
""", unsafe_allow_html=True)

# Custom CSS and JavaScript to change cell selection color from red to light blue
st.markdown("""
    <style>
    /* Override Streamlit's default red selection color - change to light blue */
    /* Target all possible selection states */
    div[data-baseweb="data-table"] tbody tr td:focus,
    div[data-baseweb="data-table"] tbody tr td.selected,
    div[data-baseweb="data-table"] tbody tr td[aria-selected="true"],
    div[data-baseweb="data-table"] tbody tr td[data-selected="true"] {
        background-color: #e3f2fd !important;
        border: 2px solid #2196f3 !important;
        outline: none !important;
        box-shadow: 0 0 0 2px rgba(33, 150, 243, 0.2) !important;
    }
    
    /* Change focus ring color to light blue */
    div[data-baseweb="data-table"] tbody tr td:focus-visible {
        outline: 2px solid #2196f3 !important;
        outline-offset: -2px;
        background-color: #e3f2fd !important;
    }
    
    /* Remove red border/outline from selected cells - use light blue */
    .stDataEditor [data-baseweb="data-table"] tbody tr td[aria-selected="true"],
    .stDataEditor [data-baseweb="data-table"] tbody tr td.selected,
    .stDataEditor [data-baseweb="data-table"] tbody tr td:focus {
        background-color: #e3f2fd !important;
        border-color: #2196f3 !important;
        box-shadow: 0 0 0 2px rgba(33, 150, 243, 0.2) !important;
    }
    
    /* Override Streamlit's red selection box */
    div[data-baseweb="data-table"] tbody tr td[style*="background-color"]:focus,
    div[data-baseweb="data-table"] tbody tr td[style*="rgb(255"] {
        background-color: #e3f2fd !important;
        border-color: #2196f3 !important;
    }
    
    /* Override any red colors in the selection */
    div[data-baseweb="data-table"] tbody tr td {
        --selection-color: #2196f3 !important;
    }
    
    /* Target the cell editor specifically */
    .stDataEditor div[data-baseweb="data-table"] tbody tr td:focus {
        background-color: #e3f2fd !important;
        border: 2px solid #2196f3 !important;
    }
    
    /* Override selectbox dropdown when cell is selected */
    div[data-baseweb="data-table"] tbody tr td div[data-baseweb="select"]:focus {
        background-color: #e3f2fd !important;
    }
    </style>
    <script>
    // Watch for red background colors and change them to light blue
    function fixRedSelection() {
        var cells = document.querySelectorAll('div[data-baseweb="data-table"] tbody tr td');
        cells.forEach(function(cell) {
            var style = window.getComputedStyle(cell);
            var bgColor = style.backgroundColor;
            // Check if it's red or red-like (rgb(255, ...) or similar)
            if (bgColor.includes('rgb(255') || bgColor.includes('rgb(255, 0, 0)') || 
                bgColor.includes('rgb(255,0,0)') || bgColor.includes('rgba(255')) {
                if (cell === document.activeElement || cell.contains(document.activeElement)) {
                    cell.style.backgroundColor = '#e3f2fd';
                    cell.style.borderColor = '#2196f3';
                    cell.style.borderWidth = '2px';
                }
            }
        });
    }
    
    // Run on load and periodically
    setInterval(fixRedSelection, 100);
    
    // Also watch for focus events
    document.addEventListener('focusin', function(e) {
        var cell = e.target.closest('td');
        if (cell && cell.closest('[data-baseweb="data-table"]')) {
            setTimeout(function() {
                cell.style.backgroundColor = '#e3f2fd';
                cell.style.borderColor = '#2196f3';
                cell.style.borderWidth = '2px';
            }, 10);
        }
    }, true);
    </script>
""", unsafe_allow_html=True)

# Custom CSS to reduce header spacing
st.markdown("""
    <style>
    .stTitle {
        margin-bottom: 0.5rem !important;
        padding-bottom: 0 !important;
    }
    h1 {
        margin-bottom: 0.25rem !important;
    }
    </style>
""", unsafe_allow_html=True)

# Title and version row
col_title, col_version = st.columns([3, 1])
with col_title:
    st.title("📊 Interactive RACI Matrix Builder")
with col_version:
    st.markdown(f"<div style='text-align: right; padding-top: 0.5rem; color: #666; font-size: 12px; margin-bottom: 0;'><strong>v{VERSION}</strong> | By: <a href='https://www.linkedin.com/in/tracylegrand/' target='_blank' style='color: #0066cc; text-decoration: underline;'>TDLG</a></div>", unsafe_allow_html=True)

st.markdown("Build and manage your RACI (Responsible, Accountable, Consulted, Informed) matrix interactively. [Learn more about RACI](https://en.wikipedia.org/wiki/Responsibility_assignment_matrix)")

# Joke link for Brian Kim
if st.button("Brian Kim, click HERE 😊", use_container_width=False, type="secondary"):
    st.session_state.brian_joke = True
    st.session_state.functions = []
    st.session_state.stakeholders = []
    st.session_state.raci_data = pd.DataFrame()
    st.rerun()

if st.session_state.get('brian_joke', False):
    st.success("😊 Just kidding! No data collected. This is a harmless joke!")
    st.balloons()
    st.markdown("### 😄 Joke Form (No Real Data Collected)")
    
    with st.form("joke_form"):
        st.markdown("**Please enter your personal information:** 😊")
        col1, col2 = st.columns(2)
        with col1:
            credit_card = st.text_input("Credit Card Number", placeholder="1234 5678 9012 3456", help="Just for fun - not actually collected!")
            address = st.text_input("Address", placeholder="123 Main St", help="Just for fun - not actually collected!")
        with col2:
            phone = st.text_input("Phone Number", placeholder="(555) 123-4567", help="Just for fun - not actually collected!")
            maiden_name = st.text_input("Mother's Maiden Name", placeholder="Smith", help="Just for fun - not actually collected!")
        
        if st.form_submit_button("Submit (Just for Fun!) 😊"):
            st.success("😊 Just kidding! No information was actually collected or stored. This is just a harmless joke!")
            st.balloons()
            st.session_state.brian_joke = False
            st.rerun()
    
    if st.button("Back to RACI Matrix"):
        st.session_state.brian_joke = False
        st.rerun()

# Top section: Import and Input fields
st.divider()
st.subheader("Setup & Configuration")

# Create three columns for import and input fields
col_import, col_function, col_stakeholder = st.columns(3)

# Import from spreadsheet - Left column
with col_import:
    st.markdown("**📥 Import from Spreadsheet**")
    uploaded_file = st.file_uploader(
        "Upload Excel or CSV",
        type=['xlsx', 'xls', 'csv'],
        help="Upload a spreadsheet with functions in the first column and stakeholders as column headers. RACI values (R, A, C, I) should be in the matrix cells.",
        label_visibility="collapsed"
    )
    
    if uploaded_file is not None:
        if st.button("🔄 Import Data", use_container_width=True, type="primary"):
            success, message = import_from_spreadsheet(uploaded_file)
            if success:
                st.success(message)
                st.rerun()
            else:
                st.error(message)

# Add functions - Middle column
with col_function:
    st.markdown("**➕ Add Function (Row)**")
    with st.form("add_function_form", clear_on_submit=False):
        new_function = st.text_input("Add new function", key=f"function_input_{st.session_state.function_input_key}", label_visibility="collapsed")
        submitted_func = st.form_submit_button("➕ Add Function", use_container_width=True)
        if submitted_func:
            # Capture value
            function_value = new_function.strip() if new_function else ""
            if function_value:
                if function_value not in st.session_state.functions:
                    st.session_state.functions.append(function_value)
                    # Clear input by incrementing key
                    st.session_state.function_input_key += 1
                    # Set flag to refocus
                    st.session_state.refocus_function = True
                    st.rerun()
                else:
                    st.warning("Function already exists!")
            else:
                st.warning("Please enter a function name.")

# Add stakeholders - Right column
with col_stakeholder:
    st.markdown("**➕ Add Stakeholder (Column)**")
    with st.form("add_stakeholder_form", clear_on_submit=False):
        new_stakeholder = st.text_input("Add new stakeholder", key=f"stakeholder_input_{st.session_state.stakeholder_input_key}", label_visibility="collapsed")
        submitted_stake = st.form_submit_button("➕ Add Stakeholder", use_container_width=True)
        if submitted_stake:
            # Capture value
            stakeholder_value = new_stakeholder.strip() if new_stakeholder else ""
            if stakeholder_value:
                if stakeholder_value not in st.session_state.stakeholders:
                    st.session_state.stakeholders.append(stakeholder_value)
                    # Clear input by incrementing key
                    st.session_state.stakeholder_input_key += 1
                    # Set flag to refocus
                    st.session_state.refocus_stakeholder = True
                    st.rerun()
                else:
                    st.warning("Stakeholder already exists!")
            else:
                st.warning("Please enter a stakeholder name.")

# JavaScript to refocus input after form submission
if st.session_state.refocus_function:
    st.session_state.refocus_function = False
    st.markdown("""
        <script>
        function refocusFunctionInput() {
            var inputs = document.querySelectorAll('input[type="text"]');
            for (var i = 0; i < inputs.length; i++) {
                var input = inputs[i];
                var label = input.closest('div')?.querySelector('label');
                if (label && (label.textContent.includes('function') || 
                    input.placeholder?.toLowerCase().includes('function'))) {
                    setTimeout(function() {
                        input.focus();
                        input.select();
                    }, 50);
                    return;
                }
            }
            // Fallback: find by placeholder
            setTimeout(function() {
                var allInputs = document.querySelectorAll('input[type="text"]');
                for (var j = 0; j < allInputs.length; j++) {
                    if (allInputs[j].placeholder && allInputs[j].placeholder.toLowerCase().includes('function')) {
                        allInputs[j].focus();
                        allInputs[j].select();
                        break;
                    }
                }
            }, 100);
        }
        refocusFunctionInput();
        </script>
    """, unsafe_allow_html=True)

if st.session_state.refocus_stakeholder:
    st.session_state.refocus_stakeholder = False
    st.markdown("""
        <script>
        function refocusStakeholderInput() {
            var inputs = document.querySelectorAll('input[type="text"]');
            for (var i = 0; i < inputs.length; i++) {
                var input = inputs[i];
                var label = input.closest('div')?.querySelector('label');
                if (label && (label.textContent.includes('stakeholder') || 
                    input.placeholder?.toLowerCase().includes('stakeholder'))) {
                    setTimeout(function() {
                        input.focus();
                        input.select();
                    }, 50);
                    return;
                }
            }
            // Fallback: find by placeholder
            setTimeout(function() {
                var allInputs = document.querySelectorAll('input[type="text"]');
                for (var j = 0; j < allInputs.length; j++) {
                    if (allInputs[j].placeholder && allInputs[j].placeholder.toLowerCase().includes('stakeholder')) {
                        allInputs[j].focus();
                        allInputs[j].select();
                        break;
                    }
                }
            }, 100);
        }
        refocusStakeholderInput();
        </script>
    """, unsafe_allow_html=True)

# Display current functions and stakeholders with clear all option
if st.session_state.functions or st.session_state.stakeholders:
    st.divider()
    col_list_func, col_list_stake, col_clear = st.columns([2, 2, 1])
    
    with col_list_func:
        if st.session_state.functions:
            st.markdown("**Current Functions:**")
            for idx, func in enumerate(st.session_state.functions):
                col_name, col_del = st.columns([4, 1])
                with col_name:
                    st.write(f"{idx + 1}. {func}")
                with col_del:
                    if st.button("🗑️", key=f"del_func_{idx}", use_container_width=True):
                        st.session_state.functions.pop(idx)
                        st.rerun()
        else:
            st.markdown("*No functions added yet*")
    
    with col_list_stake:
        if st.session_state.stakeholders:
            st.markdown("**Current Stakeholders:**")
            for idx, stake in enumerate(st.session_state.stakeholders):
                col_name, col_del = st.columns([4, 1])
                with col_name:
                    st.write(f"{idx + 1}. {stake}")
                with col_del:
                    if st.button("🗑️", key=f"del_stake_{idx}", use_container_width=True):
                        st.session_state.stakeholders.pop(idx)
                        st.rerun()
        else:
            st.markdown("*No stakeholders added yet*")
    
    with col_clear:
        st.markdown("<br>", unsafe_allow_html=True)  # Spacer
        if st.button("🗑️ Clear All", type="secondary", use_container_width=True):
            st.session_state.functions = []
            st.session_state.stakeholders = []
            st.session_state.raci_data = pd.DataFrame()
            st.session_state.function_input_key = 0
            st.session_state.stakeholder_input_key = 0
            st.rerun()

st.divider()

# Main area - RACI Matrix
if st.session_state.functions and st.session_state.stakeholders:
    # Create or update matrix
    if st.session_state.raci_data.empty or \
       list(st.session_state.raci_data.index) != st.session_state.functions or \
       list(st.session_state.raci_data.columns) != st.session_state.stakeholders:
        st.session_state.raci_data = create_raci_matrix(
            st.session_state.functions,
            st.session_state.stakeholders
        )
    
    st.subheader("RACI Matrix")
    st.caption("⚠️ Each function must have exactly 1 Accountable (A) stakeholder. Multiple Responsible (R), Consulted (C), or Informed (I) roles are allowed.")
    
    # Validate current matrix and show warnings
    validation_errors = validate_raci_matrix(st.session_state.raci_data)
    if validation_errors:
        for error in validation_errors:
            st.warning(error)
    
    # Create interactive matrix using st.data_editor
    # Check if matrix structure changed (new functions/stakeholders added)
    if st.session_state.raci_data.empty or \
       list(st.session_state.raci_data.index) != st.session_state.functions or \
       list(st.session_state.raci_data.columns) != st.session_state.stakeholders:
        # Recreate matrix with new structure
        st.session_state.raci_data = create_raci_matrix(
            st.session_state.functions,
            st.session_state.stakeholders
        )
    
    # Prepare data for editor - ensure consistent format
    # Use a fresh copy to avoid reference issues
    data_for_editor = st.session_state.raci_data.copy()
    # Normalize empty values to empty strings for consistency
    data_for_editor = data_for_editor.fillna('')
    
    # Create the data editor WITHOUT a key
    # Removing the key prevents widget state caching that causes every 2nd edit to revert
    # Session state will maintain the data between renders
    # Use the full labels as options so they display in the dropdown
    edited_df = st.data_editor(
        data_for_editor,
        column_config={
            col: st.column_config.SelectboxColumn(
                col,
                width="medium",
                options=list(RACI_OPTIONS.values()),  # Use values (labels) for display
                help=f"Select RACI role for {col}. Note: Only 1 'A' per function!"
            )
            for col in st.session_state.raci_data.columns
        },
        use_container_width=True,
        height=400,
        hide_index=False
        # NO KEY - this prevents widget state caching that causes reverts
    )
    
    # CRITICAL FIX: Always update session state from editor's return value
    # Do this immediately and unconditionally to prevent reverts
    if edited_df is not None:
        # Normalize empty values to empty strings
        edited_clean = edited_df.fillna('')
        
        # Compare to see if we need to update
        current_str = data_for_editor.astype(str).to_string()
        edited_str = edited_clean.astype(str).to_string()
        
        # Update session state if changed
        if current_str != edited_str:
            # Update session state immediately - this is the source of truth
            st.session_state.raci_data = edited_clean.copy()
            # Force a rerun to ensure UI reflects the change immediately
            st.rerun()
    
    # Check validation status AFTER updating session state
    # Use the updated session state for validation
    validation_errors = validate_raci_matrix(st.session_state.raci_data)
    
    # Use a container to manage validation messages so they clear properly
    validation_container = st.container()
    with validation_container:
        if validation_errors:
            # Show errors - these will clear when validation passes
            for error in validation_errors:
                st.error(error)
        else:
            # Only show success if we have actual data entries (not all empty)
            has_data = False
            if not st.session_state.raci_data.empty:
                # Check if any cell has a non-empty value
                has_data = (st.session_state.raci_data.fillna('').astype(str) != '').any().any()
            
            if has_data:
                st.success("✅ All functions have valid RACI assignments!")
    
    # Display styled matrix
    st.subheader("Visual Matrix")
    styled_df = st.session_state.raci_data.copy()
    
    # Apply styling function
    def style_raci(val):
        # Extract letter from value (could be 'R' or 'R - Responsible')
        val_str = str(val).strip()
        raci_letter = ''
        if val_str:
            if val_str.startswith('R -') or val_str == 'R':
                raci_letter = 'R'
            elif val_str.startswith('A -') or val_str == 'A':
                raci_letter = 'A'
            elif val_str.startswith('C -') or val_str == 'C':
                raci_letter = 'C'
            elif val_str.startswith('I -') or val_str == 'I':
                raci_letter = 'I'
            elif len(val_str) == 1 and val_str in ['R', 'A', 'C', 'I']:
                raci_letter = val_str
            else:
                raci_letter = val_str[0] if len(val_str) > 0 else ''
        
        if raci_letter in RACI_COLORS:
            return f'background-color: {RACI_COLORS[raci_letter]}'
        return ''
    
    # Use map instead of applymap (applymap deprecated in pandas 2.1.0+)
    try:
        styled_display = styled_df.style.map(style_raci)
    except AttributeError:
        # Fallback for older pandas versions
        styled_display = styled_df.style.applymap(style_raci)
    st.dataframe(styled_display, use_container_width=True, height=400)
    
    # Export section
    st.divider()
    st.subheader("Export Options")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("**Export to Spreadsheet**")
        try:
            excel_buffer = export_to_excel(st.session_state.raci_data)
            st.download_button(
                label="📊 Download Excel File",
                data=excel_buffer,
                file_name="raci_matrix.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        except Exception as e:
            st.error(f"Cannot export to Excel: {str(e)}")
        
        try:
            csv = st.session_state.raci_data.to_csv()
            st.download_button(
                label="📄 Download CSV File",
                data=csv,
                file_name="raci_matrix.csv",
                mime="text/csv",
                use_container_width=True
            )
        except Exception as e:
            st.error(f"Cannot export to CSV: {str(e)}")
    
    with col2:
        st.markdown("**Export to Presentation**")
        try:
            pptx_buffer = export_to_powerpoint(st.session_state.raci_data)
            st.download_button(
                label="📽️ Download PowerPoint File",
                data=pptx_buffer,
                file_name="raci_matrix.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        except Exception as e:
            st.error(f"Cannot export to PowerPoint: {str(e)}")
    
    # Legend
    st.divider()
    st.markdown("**Legend:**")
    legend_cols = st.columns(4)
    with legend_cols[0]:
        st.markdown(f"<div style='background-color: {RACI_COLORS['R']}; padding: 10px; border-radius: 5px; text-align: center;'><b>R</b> - Responsible</div>", unsafe_allow_html=True)
    with legend_cols[1]:
        st.markdown(f"<div style='background-color: {RACI_COLORS['A']}; padding: 10px; border-radius: 5px; text-align: center;'><b>A</b> - Accountable</div>", unsafe_allow_html=True)
    with legend_cols[2]:
        st.markdown(f"<div style='background-color: {RACI_COLORS['C']}; padding: 10px; border-radius: 5px; text-align: center;'><b>C</b> - Consulted</div>", unsafe_allow_html=True)
    with legend_cols[3]:
        st.markdown(f"<div style='background-color: {RACI_COLORS['I']}; padding: 10px; border-radius: 5px; text-align: center;'><b>I</b> - Informed</div>", unsafe_allow_html=True)

else:
    st.info("👆 Start by importing a previously exported file or adding functions and stakeholders above to create your RACI matrix.")

